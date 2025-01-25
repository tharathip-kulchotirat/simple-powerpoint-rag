from pathlib import Path
from pptx import Presentation
from loguru import logger
from src.ppt_chatbot.utils.file_handling import secure_tempfile

class PPTParser:
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def parse_pptx(self, file_path: Path) -> list[str]:
        """Convert PPTX file to structured text chunks"""
        try:
            presentation = Presentation(file_path)
            return self._process_slides(presentation)
        except Exception as e:
            logger.error(f"PPT parsing failed: {str(e)}")
            raise

    def _process_slides(self, presentation) -> list[str]:
        """Process individual slides"""
        chunks = []
        for i, slide in enumerate(presentation.slides):
            slide_content = [
                f"# Slide {i+1}",
                self._extract_text(slide)
            ]
            chunks.extend(self._chunk_content('\n\n'.join(slide_content)))
        return chunks

    def _extract_text(self, slide) -> str:
        """Extract text from slide elements"""
        return '\n'.join(
            shape.text.strip() 
            for shape in slide.shapes 
            if hasattr(shape, "text")
        )

    def _chunk_content(self, content: str) -> list[str]:
        """Split content into manageable chunks"""
        words = content.split()
        return [' '.join(words[i:i+self.chunk_size]) 
                for i in range(0, len(words), self.chunk_size - self.chunk_overlap)]